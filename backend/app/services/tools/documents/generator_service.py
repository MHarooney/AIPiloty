"""Document Generator Service — creates PDF, XLSX, DOCX, PPTX files safely."""

from __future__ import annotations

import json
import logging
import os
import re
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

MAX_SECTIONS = 50
MAX_STRING_FIELD_BYTES = 100_000
MAX_PAYLOAD_BYTES = 2_000_000
MAX_FILENAME_LEN = 200
GENERATED_DIR_NAME = "generated"

_SAFE_FILENAME_RE = re.compile(r"[^a-zA-Z0-9_\-. ]")


def _sanitize_filename(name: str) -> str:
    name = name.replace("\x00", "").replace("/", "_").replace("\\", "_")
    name = _SAFE_FILENAME_RE.sub("_", name)
    name = re.sub(r"_+", "_", name).strip("_ ")
    if not name:
        name = f"document_{uuid.uuid4().hex[:8]}"
    return name[:MAX_FILENAME_LEN]


def _auto_filename(title: str, ext: str) -> str:
    base = _sanitize_filename(title) if title else f"document_{uuid.uuid4().hex[:8]}"
    if not base.lower().endswith(f".{ext}"):
        base = f"{base}.{ext}"
    return base


def _validate_string_field(value: Any, field_name: str) -> str:
    s = str(value) if value is not None else ""
    if len(s.encode("utf-8", errors="replace")) > MAX_STRING_FIELD_BYTES:
        raise ValueError(f"Field '{field_name}' exceeds {MAX_STRING_FIELD_BYTES // 1000} KB limit")
    return s


def _xml_escape(s: str) -> str:
    """Escape text for ReportLab Paragraph (HTML-like)."""
    return (
        s.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


class DocumentGeneratorService:
    """Generate Office/PDF documents in a workspace-safe directory."""

    def __init__(self, workspace_root: str):
        self._root = str(Path(workspace_root).resolve())
        self._generated_dir = os.path.join(self._root, GENERATED_DIR_NAME)
        os.makedirs(self._generated_dir, exist_ok=True)

    def _safe_output_path(self, filename: str) -> str:
        sanitized = _sanitize_filename(filename)
        resolved = str(Path(os.path.join(self._generated_dir, sanitized)).resolve())
        if not resolved.startswith(self._generated_dir):
            raise ValueError(f"Filename resolves outside generated directory")
        return resolved

    def _make_result(self, resolved_path: str) -> Dict[str, Any]:
        size = os.path.getsize(resolved_path)
        rel = os.path.relpath(resolved_path, self._root)
        ext = os.path.splitext(resolved_path)[1].lower()
        mime_map = {
            ".pdf": "application/pdf",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".png": "image/png",
        }
        return {"success": True, "relative_path": rel, "bytes": size, "mime_type": mime_map.get(ext, "application/octet-stream")}

    @staticmethod
    def _parse_json_param(raw: str, param_name: str) -> Any:
        if not isinstance(raw, str):
            return raw
        try:
            return json.loads(raw)
        except (json.JSONDecodeError, TypeError) as e:
            raise ValueError(f"Parameter '{param_name}' is not valid JSON: {e}")

    @staticmethod
    def _check_payload_size(data: Any) -> None:
        if len(json.dumps(data, default=str).encode()) > MAX_PAYLOAD_BYTES:
            raise ValueError(f"Payload exceeds {MAX_PAYLOAD_BYTES // 1_000_000} MB limit")

    # ── PDF ──────────────────────────────────────────────────────

    async def generate_pdf(self, title: str, sections: str, filename: Optional[str] = None) -> Dict[str, Any]:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
        from reportlab.lib import colors

        title = _validate_string_field(title, "title")
        parsed = self._parse_json_param(sections, "sections")
        if not isinstance(parsed, list):
            raise ValueError("'sections' must be a JSON array")
        if len(parsed) > MAX_SECTIONS:
            raise ValueError(f"Too many sections ({len(parsed)}); max is {MAX_SECTIONS}")
        self._check_payload_size({"title": title, "sections": parsed})

        fname = _auto_filename(filename or title, "pdf")
        out_path = self._safe_output_path(fname)
        styles = getSampleStyleSheet()
        code_style = ParagraphStyle(
            "CodeBlock",
            parent=styles["BodyText"],
            fontName="Courier",
            fontSize=9,
            leading=11,
            leftIndent=12,
            rightIndent=12,
            spaceBefore=4,
            spaceAfter=8,
            backColor=colors.HexColor("#f4f4f4"),
        )
        story = []
        story.append(Paragraph(_xml_escape(title), styles["Title"]))
        story.append(Spacer(1, 0.5 * cm))

        for sec in parsed:
            heading = _validate_string_field(sec.get("heading", ""), "section.heading")
            body = _validate_string_field(sec.get("body", ""), "section.body")
            if heading:
                story.append(Paragraph(_xml_escape(heading), styles["Heading2"]))
                story.append(Spacer(1, 0.2 * cm))
            if body:
                for para in body.split("\n\n"):
                    p = para.strip()
                    if p:
                        story.append(Paragraph(_xml_escape(p), styles["BodyText"]))
                story.append(Spacer(1, 0.2 * cm))
            bullets = sec.get("bullets")
            if isinstance(bullets, list) and bullets:
                for item in bullets:
                    story.append(
                        Paragraph(f"• {_xml_escape(_validate_string_field(item, 'bullet'))}", styles["BodyText"])
                    )
                story.append(Spacer(1, 0.2 * cm))
            code = sec.get("code")
            if isinstance(code, str) and code.strip():
                code = _validate_string_field(code, "section.code")
                code_html = _xml_escape(code).replace("\n", "<br/>")
                story.append(Paragraph(code_html, code_style))
                story.append(Spacer(1, 0.2 * cm))
            diagram = sec.get("diagram")
            if isinstance(diagram, str) and diagram.strip():
                d = _validate_string_field(diagram, "section.diagram")
                d_html = _xml_escape(d).replace("\n", "<br/>")
                story.append(Paragraph(f"<i>{d_html}</i>", styles["BodyText"]))
                story.append(Spacer(1, 0.2 * cm))
            table_data = sec.get("table")
            if table_data and isinstance(table_data, list) and len(table_data) > 0:
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4472C4")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                ]))
                story.append(t)
                story.append(Spacer(1, 0.3 * cm))

        doc = SimpleDocTemplate(out_path, pagesize=A4)
        doc.build(story)
        return self._make_result(out_path)

    # ── XLSX ─────────────────────────────────────────────────────

    async def generate_xlsx(self, title: str, sheets: str, filename: Optional[str] = None) -> Dict[str, Any]:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill

        title = _validate_string_field(title, "title")
        parsed = self._parse_json_param(sheets, "sheets")
        if not isinstance(parsed, list):
            raise ValueError("'sheets' must be a JSON array")
        self._check_payload_size({"title": title, "sheets": parsed})

        fname = _auto_filename(filename or title, "xlsx")
        out_path = self._safe_output_path(fname)
        wb = Workbook()
        default_ws = wb.active

        for idx, sheet_def in enumerate(parsed):
            name = str(sheet_def.get("name", f"Sheet{idx + 1}"))[:31]
            headers = sheet_def.get("headers", [])
            rows = sheet_def.get("rows", [])
            if idx == 0 and default_ws is not None:
                ws = default_ws
                ws.title = name
            else:
                ws = wb.create_sheet(title=name)
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF")
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=str(h))
                cell.fill = header_fill
                cell.font = header_font
            start_row = 2 if headers else 1
            for r_idx, row in enumerate(rows):
                if not isinstance(row, (list, tuple)):
                    row = [row]
                for c_idx, val in enumerate(row, 1):
                    ws.cell(row=start_row + r_idx, column=c_idx, value=val)
            for col in ws.columns:
                max_len = 0
                for cell in col:
                    try:
                        if cell.value:
                            max_len = max(max_len, len(str(cell.value)))
                    except Exception:
                        pass
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 50)

        wb.save(out_path)
        return self._make_result(out_path)

    # ── DOCX ─────────────────────────────────────────────────────

    async def generate_docx(self, title: str, sections: str, filename: Optional[str] = None) -> Dict[str, Any]:
        from docx import Document
        from docx.shared import Pt

        title = _validate_string_field(title, "title")
        parsed = self._parse_json_param(sections, "sections")
        if not isinstance(parsed, list):
            raise ValueError("'sections' must be a JSON array")
        self._check_payload_size({"title": title, "sections": parsed})

        fname = _auto_filename(filename or title, "docx")
        out_path = self._safe_output_path(fname)
        doc = Document()
        doc.add_heading(title, level=0)

        for sec in parsed:
            heading = _validate_string_field(sec.get("heading", ""), "section.heading")
            body = _validate_string_field(sec.get("body", ""), "section.body")
            level = sec.get("level", 1)
            if not isinstance(level, int) or level < 1 or level > 4:
                level = 1
            if heading:
                doc.add_heading(heading, level=level)
            if body:
                for para_text in body.split("\n"):
                    if para_text.strip():
                        p = doc.add_paragraph(para_text.strip())
                        p.style.font.size = Pt(11)
            bullets = sec.get("bullets")
            if bullets and isinstance(bullets, list):
                for item in bullets:
                    doc.add_paragraph(str(item), style="List Bullet")

        doc.save(out_path)
        return self._make_result(out_path)

    # ── PPTX ─────────────────────────────────────────────────────

    async def generate_pptx(self, title: str, slides: str, filename: Optional[str] = None) -> Dict[str, Any]:
        from pptx import Presentation

        title = _validate_string_field(title, "title")
        parsed = self._parse_json_param(slides, "slides")
        if not isinstance(parsed, list):
            raise ValueError("'slides' must be a JSON array")
        self._check_payload_size({"title": title, "slides": parsed})

        fname = _auto_filename(filename or title, "pptx")
        out_path = self._safe_output_path(fname)
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1]:
            slide.placeholders[1].text = "Generated by AIPiloty"

        content_layout = prs.slide_layouts[1]
        for slide_def in parsed:
            s_title = _validate_string_field(slide_def.get("title", ""), "slide.title")
            s_content = _validate_string_field(slide_def.get("content", ""), "slide.content")
            sl = prs.slides.add_slide(content_layout)
            sl.shapes.title.text = s_title
            body = sl.placeholders[1]
            tf = body.text_frame
            tf.text = ""
            for i, line in enumerate(s_content.split("\n")):
                if i == 0:
                    tf.text = line.strip()
                else:
                    p = tf.add_paragraph()
                    p.text = line.strip()

        prs.save(out_path)
        return self._make_result(out_path)

    # ── Image (external API) ─────────────────────────────────────

    async def generate_image(self, prompt: str, width: int = 512, height: int = 512, filename: Optional[str] = None) -> Dict[str, Any]:
        import httpx

        api_url = os.environ.get("IMAGE_GEN_API_URL", "").strip()
        if not api_url:
            return {"success": False, "error": "Image generation not configured. Set IMAGE_GEN_API_URL."}

        prompt = _validate_string_field(prompt, "prompt")
        width = max(64, min(int(width), 2048))
        height = max(64, min(int(height), 2048))
        fname = _auto_filename(filename or f"image_{uuid.uuid4().hex[:8]}", "png")
        out_path = self._safe_output_path(fname)

        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                resp = await client.post(f"{api_url.rstrip('/')}/generate", json={"prompt": prompt, "width": width, "height": height})
                resp.raise_for_status()
                content_type = resp.headers.get("content-type", "")
                if "image" in content_type:
                    with open(out_path, "wb") as f:
                        f.write(resp.content)
                elif "json" in content_type:
                    import base64
                    data = resp.json()
                    img_b64 = data.get("image") or data.get("images", [None])[0]
                    if not img_b64:
                        return {"success": False, "error": "API returned no image data"}
                    with open(out_path, "wb") as f:
                        f.write(base64.b64decode(img_b64))
                else:
                    return {"success": False, "error": f"Unexpected content type: {content_type}"}
        except httpx.HTTPStatusError as e:
            return {"success": False, "error": f"Image API {e.response.status_code}"}
        except httpx.RequestError as e:
            return {"success": False, "error": f"Image API connection error: {e}"}

        return self._make_result(out_path)
