"""PPTX text extraction using python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def extract(file_path: Path) -> str:
    """Extract text from all slides in a PPTX file."""
    prs = Presentation(str(file_path))
    parts = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            parts.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
    return "\n\n".join(parts)
